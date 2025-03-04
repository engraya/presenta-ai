from fastapi import FastAPI, BackgroundTasks
from pydantic import BaseModel, Field 
from typing import Optional
from openai import OpenAI
import os
from dotenv import load_dotenv  
from pptx import Presentation
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware 
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# Load environment variables
load_dotenv()

# Get the API key from environment
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("OPENAI_API_KEY is not set in the environment or .env file!")

# Initialize OpenAI Client
client = OpenAI(api_key=api_key)

app = FastAPI()

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)



class PPTRequest(BaseModel):
    topic: str
    num_slides: Optional[int] = Field(5, ge=1, le=20)
    layout_preference: Optional[str] = "Varied"

def generate_slide_content(topic: str, num_slides: int):
    prompt = f"""
    Generate content for a {num_slides}-slide PowerPoint presentation about "{topic}".
    Each slide should have:
    - A clear slide title
    - 3-5 key bullet points explaining the topic
    - If relevant, suggest an image placeholder with a description (e.g., "Insert Image: Growth of AI in Healthcare").
    
    Ensure the content is evenly distributed across {num_slides} slides.

    Format the response as:
    
    Slide 1:
    Title: [Slide Title]
    - [Bullet Point 1]
    - [Bullet Point 2]
    - [Bullet Point 3]
    Image: [Optional Image Placeholder]
    
    Slide 2:
    Title: [Slide Title]
    - [Bullet Point 1]
    - [Bullet Point 2]
    - [Bullet Point 3]
    Image: [Optional Image Placeholder]

    (Continue this format for all {num_slides} slides)
    """

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a PowerPoint expert."},
            {"role": "user", "content": prompt}
        ]
    )
    
    response_message = response.choices[0].message.content.strip()
    
    slides = response_message.split("\n\n")
    structured_slides = []

    for slide in slides:
        lines = slide.split("\n")
        title = lines[0].replace("Title: ", "").strip()
        bullets = [line.strip() for line in lines[1:] if line.startswith("- ")]
        image = next((line.replace("Image: ", "").strip() for line in lines if line.startswith("Image: ")), None)

        structured_slides.append({"title": title, "bullets": bullets, "image": image})
    
    return structured_slides


def create_pptx(topic: str, slides_data, layout_preference: str = "default", filename="presentation.pptx"):
        # Load a predefined PowerPoint template
    prs = Presentation("template.pptx")  # Ensure correct template path
    
    # Title Slide (First slide should always have the title)
    title_slide_layout = prs.slide_layouts[0]  
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = topic
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)  
    slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if len(slide.placeholders) > 1:  
        slide.placeholders[1].text = "AI-Generated Presentation"
        slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(24)
        slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Define Layout Mapping
    layout_mapping = {
        "Varied": [1, 5, 3, 2, 8],  
        "Text-Heavy": [1, 1, 1, 1, 1],  
        "Image-Focused": [8, 8, 8, 8, 8],  
        "default": [1, 5, 3, 2, 8],  
    }

    for i, slide_content in enumerate(slides_data):
        title_text = slide_content.get("title", "").strip()
        bullet_points = slide_content.get("bullets", [])
        image_path = slide_content.get("image", None)

        # Skip blank slides (no title, bullets, or image)
        if not title_text and not bullet_points and not image_path:
            continue

        # Choose layout dynamically
        chosen_layout = layout_mapping.get(layout_preference, layout_mapping["default"])
        slide_layout_index = chosen_layout[i % len(chosen_layout)]
        slide_layout = prs.slide_layouts[slide_layout_index]

        slide = prs.slides.add_slide(slide_layout)

        # Add Title at the Top
        if title_text:
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title_text[:50]  
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.size = Pt(32 if len(title_text) <= 30 else 26)
                title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Handle Text-Based Layouts (1, 2, 5) - Adding Bullets
        if slide_layout_index in [1, 2, 5] and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]  
            content_placeholder.text = ""  # Clear existing text

            for bullet in bullet_points:
                p = content_placeholder.text_frame.add_paragraph()
                p.text = bullet
                p.space_after = Pt(10)
                p.font.size = Pt(22 if len(bullet_points) <= 5 else 18)  
                p.level = 0  # Ensure it's formatted as a bullet point

        # Handle Two-Column Layout (3)
        elif slide_layout_index == 3 and len(slide.placeholders) > 1:
            left_text = slide.placeholders[0]
            right_text = slide.placeholders[1]

            left_text.text = title_text
            for bullet in bullet_points:
                p = right_text.text_frame.add_paragraph()
                p.text = bullet
                p.space_after = Pt(10)

        # Handle Image Layout (8) - Ensure Title is at the Top
        elif slide_layout_index == 8:
            left = Inches(1)
            top = Inches(2.5)  # Ensuring it's placed below the title

            if image_path:
                try:
                    left = Inches(1)
                    top = Inches(3.5)  # Ensure it is below the title
                    width = Inches(7.5)
                    height = Inches(4.5)
                    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                except Exception:
                    if slide_layout in [prs.slide_layouts[5], prs.slide_layouts[8]]:  # Only for image slides
                        text_placeholder = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(7.5), Inches(1))
                        text_placeholder.text = f"Insert Image: {topic.replace(' ', '_')}"
                        text_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                        text_placeholder.text_frame.paragraphs[0].font.size = Pt(16)
                        text_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            elif slide_layout == prs.slide_layouts[8]:  # If it's an image slide but no image provided
                text_placeholder = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(7.5), Inches(1))
                text_placeholder.text = f"Insert Image: {topic.replace(' ', '_')}"
                text_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
                text_placeholder.text_frame.paragraphs[0].font.size = Pt(16)
                text_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(filename)
    return filename


@app.get("/")
async def welcome():
    return {"message": "Welcome to Presenta-AI! Use /generate_ppt to create a PowerPoint presentation."}

@app.post("/generate_ppt")
async def generate_ppt(request: PPTRequest, background_tasks: BackgroundTasks):
    slides_data = generate_slide_content(request.topic, request.num_slides)
    filename = f"{request.topic.replace(' ', '_')}.pptx"
    background_tasks.add_task(create_pptx, request.topic, slides_data, request.layout_preference, filename)
    return {"message": "Presentation generated successfully!.", "filename": filename}

@app.get("/download_ppt/{filename}")
async def download_ppt(filename: str):
    file_path = f"./{filename}"
    if os.path.exists(file_path):
        return FileResponse(file_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)
    return {"error": "File not found"}
